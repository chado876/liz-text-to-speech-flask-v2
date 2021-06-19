import speechUtil as speechUtil
import os
import glob
import PyPDF2
import docx2txt
from pptx import Presentation
from nlp.lex.lexical_analyzer import LexicalAnalyzer
from nltk import Tree, tree
from nltk.draw.util import CanvasFrame
from nltk.draw import TreeWidget
import os
from PyPDF2 import PdfFileMerger

def readFromFile(fileName, treeFileName):
    filePath = "uploads/" + fileName
    filetype = fileName.split('.')[-1]
    fileNameWithoutExt = ( fileName.rsplit( ".", 1 )[ 0 ] )
    success = True
    text = ''
    
    try:
        if filetype == 'txt':
            if is_file_empty(filePath):
                success = False
                raise Exception("File string is empty!")
            else:
                with open(filePath,"r") as f:
                    text = f.read()
                    speechUtil.synthesize_and_save_to_file(text, fileNameWithoutExt)
        elif filetype == 'pdf':
            pdfFileObj = open(filePath, "rb")
            pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
            totalPages = pdfReader.numPages
            for i in range(totalPages):
                page = pdfReader.getPage(i)
                text += page.extractText()
            speechUtil.synthesize_and_save_to_file(text, fileNameWithoutExt)
        elif filetype == 'docx':
            text = docx2txt.process(filePath)
            speechUtil.synthesize_and_save_to_file(text, fileNameWithoutExt)
        elif filetype == 'pptx':
            pptOutput = ''
            ppt = Presentation(filePath)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            speechUtil.synthesize_and_save_to_file(text, fileNameWithoutExt)
    except:
        print("An error occurred!")
    return text
        



def clean_up_files():
    files = glob.glob('output/*')
    for f in files:
        os.remove(f)
    files2 = glob.glob('uploads/*')
    for f in files2:
        os.remove(f)
    files3 = glob.glob('parse_output/*')
    for f in files3:
        os.remove(f)

def is_file_empty(filePath):
    with open(filePath, 'r') as fileReader:
        first_char = fileReader.read(1)
        if not first_char:
            return True
    return False


def generate_tree_pdf(extractions, treeFileName):
    merger = PdfFileMerger()

    for x,extraction in enumerate(extractions):
        print(str(extraction))
        cf = CanvasFrame()
        t = Tree.fromstring(str(extraction))
        tc = TreeWidget(cf.canvas(),t)
        tc['node_font'] = 'arial 14 bold'
        tc['leaf_font'] = 'arial 14'
        tc['node_color'] = '#005990'
        tc['leaf_color'] = '#3F8F57'
        tc['line_color'] = '#175252'
        cf.add_widget(tc,10,10) # (10,10) offsets
        ps_file_name = 'tree'+"-"+str(x)+".ps"
        pdf_file_name = 'tree'+"-"+str(x)+".pdf"
        cf.print_to_file('parse_output/' + ps_file_name)
        cf.destroy()
        os.system('magick convert parse_output/' + ps_file_name + ' parse_output/' + pdf_file_name)

        merger.append('parse_output/' + pdf_file_name)
    
    merger.write('parse_output/' + treeFileName)
    merger.close()
